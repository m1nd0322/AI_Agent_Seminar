#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI Agent 도입 실패 내용을 PowerPoint 프레젠테이션에 추가하는 스크립트
6페이지(인덱스 5) 이후에 슬라이드를 삽입
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

def add_content_slide(prs, title, content_items, insert_index=None):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout

    if insert_index is not None:
        # 특정 위치에 슬라이드 삽입
        import copy
        slide = prs.slides.add_slide(slide_layout)

        # 슬라이드를 원하는 위치로 이동
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)
        xml_slides.remove(slides_list[-1])
        xml_slides.insert(insert_index, slides_list[-1])
    else:
        slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # 들여쓰기 레벨 설정
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0

        # 폰트 크기 조정
        if len(item if isinstance(item, str) else item[0]) > 80:
            p.font.size = Pt(12)
        else:
            p.font.size = Pt(14)

    return slide

def insert_slides_after_index(prs, slides_data, insert_after_index):
    """특정 인덱스 이후에 여러 슬라이드를 삽입"""
    xml_slides = prs.slides._sldIdLst

    # 먼저 모든 슬라이드를 끝에 추가
    new_slides = []
    for title, content_items in slides_data:
        slide = add_content_slide(prs, title, content_items)
        new_slides.append(slide)

    # XML에서 슬라이드 순서 재배열
    slides_list = list(xml_slides)
    num_new_slides = len(new_slides)

    # 새로 추가된 슬라이드들을 끝에서 가져와서 원하는 위치로 이동
    for i in range(num_new_slides):
        slide_element = slides_list[-(num_new_slides - i)]
        xml_slides.remove(slide_element)
        xml_slides.insert(insert_after_index + 1 + i, slide_element)

def main():
    # 기존 PowerPoint 파일 열기
    prs = Presentation('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')

    print(f"기존 슬라이드 수: {len(prs.slides)}")

    # 6페이지(인덱스 5) 이후에 삽입할 슬라이드 데이터
    slides_data = [
        # 슬라이드 1: Executive Summary
        ("AI Agent 도입 실패: Executive Summary", [
            "최근 2년간 기업들의 생성형 AI(GenAI) 대규모 도입",
            ("실질적 생산성 향상과 사업 가치 창출은 극소수(약 5%)에 한정", 1),
            "",
            "도입 실패의 주요 원인:",
            ("기술력 부족이 아닌 내부 정치, 레거시 시스템, 조직 설계 부재", 1),
            "",
            "2025년 '완전 자율 AI 에이전트의 해' 과대광고:",
            ("프로덕션급 완전 자율 에이전트는 여전히 실현 불가능", 1),
            "",
            "기술적 현실:",
            ("완전 자율 에이전트는 에러 누적과 비용 문제로 시기상조", 1),
            ("'Human-in-the-loop' 기반 도구 설계 필수", 1)
        ]),

        # 슬라이드 2: 브린욜프슨 생산성 역설
        ("브린욜프슨 생산성 역설 2.0", [
            "역사적 사례: 1990년대 PC·IT 대규모 투자",
            ("생산성 지표는 10~15년간 거의 상승하지 않음", 1),
            "",
            "현재 상황: 'High Adoption, Low Transformation'",
            ("GenAI 도입률은 매우 높음", 1),
            ("경제 전체 생산성 지표 및 기업 내부 구조 변화는 미미", 1),
            ("단순 도입 단계에 머물러 있음", 1),
            "",
            "핵심 원인:",
            ("AI는 단독 기술로는 생산성을 끌어올리지 못함", 1),
            ("조직 재설계 + 업무 프로세스 재구성 + 직원 스킬 변화", 1),
            ("보완재(Complementary Investment)와 함께 할 때만 효과 발생", 1)
        ]),

        # 슬라이드 3: MIT 연구 - GenAI Divide
        ("MIT 연구: GenAI Divide (95% vs 5%)", [
            "성공 기업 (약 5%):",
            ("실질 ROI 발생, 매출·생산성 개선", 1),
            ("학습·기억·적응이 가능한 에이전트형 시스템 구축", 1),
            ("실제 업무 흐름에 깊게 통합", 1),
            "",
            "실패 기업 (약 95%):",
            ("투자만 하고 성과 없음", 1),
            ("기술 도입에 그침", 1),
            "",
            "의미 있는 변화:",
            ("테크·미디어·금융 등 정보 중심 산업에 집중", 1)
        ]),

        # 슬라이드 4: 대기업 실패 이유 1
        ("대기업이 AI를 실패하는 구조적 이유 (1/2)", [
            "1. 학습 격차 (Learning Gap)",
            ("대부분의 엔터프라이즈 AI는 정적 도구", 1),
            ("사용자 피드백 학습 불가", 1),
            ("단순 작업: ChatGPT / 복잡 작업: 여전히 사람", 1),
            "",
            "2. 파일럿-프로덕션 단절",
            ("POC는 수백 개, 전사 배포 성공률 5% 미만", 1),
            ("실제 업무 프로세스와 불일치", 1),
            ("인프라·규제 문제가 아닌 워크플로·조직 설계 문제", 1),
            "",
            "3. 인프라 문제",
            ("내부 IT 시스템의 노후화 및 사일로화", 1),
            ("데이터 통합 불가", 1)
        ]),

        # 슬라이드 5: 대기업 실패 이유 2
        ("대기업이 AI를 실패하는 구조적 이유 (2/2)", [
            "4. 조직-정치적 장벽",
            ("여러 조직 간 이해관계 충돌", 1),
            ("부서 간 협력 부재", 1),
            "",
            "5. 인적 자원",
            ("내부 엔지니어링 팀의 AI 회의론 및 변화 거부", 1),
            ("코드 생성 도구조차 사용하지 않는 경향", 1),
            "",
            "6. 그림자 AI (Shadow AI) 확산",
            ("공식 사내 도구가 불편하고 규제가 강함", 1),
            ("직원들이 개인 계정 GPT·Claude·Gemini로 실무 처리", 1),
            ("숨은 수요가 매우 크다는 강력한 신호", 1),
            "",
            "시사점: 검증된 외부 스타트업 솔루션 도입이 성공 확률 높음"
        ]),

        # 슬라이드 6: 2025년 자율 AI 에이전트 과대광고
        ("2025년 자율 AI 에이전트: 과대광고 vs 현실", [
            "과대광고: 완전 자율 다단계 워크플로우 가능",
            "현실: 불가능",
            "",
            "과대광고: '거의 안 틀린다'",
            "현실: 에러율 문제",
            ("각 단계 성공률 95%라도 20단계면 성공률 36.8%", 1),
            ("프로덕션 요구치 99.9%+ 불충족", 1),
            "",
            "과대광고: 비용 무시",
            "현실: 토큰 비용 폭증",
            ("대화 길어질수록 비용 제곱 증가", 1),
            ("100회 길이 대화: 인당 50~100달러", 1),
            ("대량 사용자 시 경제성 붕괴", 1)
        ]),

        # 슬라이드 7: 실제 작동하는 에이전트
        ("실제 작동하는 에이전트 시스템의 현실", [
            "성공 패턴:",
            ("제한된 도메인 + 단계별 사람 검증 필수", 1),
            ("대화형 만능 에이전트 불가능", 1),
            "",
            "실제 작동하는 시스템의 70%는 전통 엔지니어링:",
            ("대량 정보 요약 (상태 추상화)", 1),
            ("툴 실패 시 정보 전달·복구 로직", 1),
            ("컨텍스트 오염 방지", 1),
            ("구조화된 피드백을 통한 의사결정 가능 상태 제공", 1),
            "",
            "시사점:",
            ("'AI가 전체 스택을 완전 자율 운영'은 현실의 벽과 충돌", 1)
        ]),

        # 슬라이드 8: 종합 결론 1
        ("종합 결론 및 향후 전략 (1/2)", [
            "1) 조직 및 인력 전략 (HR & Culture)",
            "",
            "Shadow AI 양성화:",
            ("직원들이 개별적으로 사용하는 AI 활용 패턴 분석", 1),
            ("공식 워크플로에 반영하는 Bottom-up 혁신 유도", 1),
            "",
            "2) 개발 및 도입 전략 (Dev & Tech)",
            "",
            "Human-in-the-loop 필수화:",
            ("완전 자율 에이전트 환상 버리기", 1),
            ("중요 의사결정 단계에 인간 개입하는 반자동화 시스템", 1)
        ]),

        # 슬라이드 9: 종합 결론 2
        ("종합 결론 및 향후 전략 (2/2)", [
            "명확한 도구(Tool) 중심 설계:",
            ("만능 대화형 챗봇보다는", 1),
            ("특정 업무를 완벽히 수행하는 '스마트 도구'들의 집합", 1),
            "",
            "전통 엔지니어링 강화:",
            ("AI 모델 튜닝보다", 1),
            ("AI 실패 시 복구(Fallback)와 시스템 안정성 유지에 투자", 1),
            ("백엔드 엔지니어링 역량 확보", 1),
            "",
            "3) 의사결정 방향",
            "",
            "Buy over Build:",
            ("내부 정치와 레거시 문제로 자체 개발 난항 시", 1),
            ("검증된 버티컬 AI 솔루션 신속 도입", 1),
            ("'학습 격차' 해소", 1)
        ]),

        # 슬라이드 10: 최종 메시지
        ("AI Agent 도입 성공의 핵심", [
            "생성형 AI의 진짜 승부는 이제부터 시작",
            "",
            "성공의 열쇠:",
            ("기술이 아니라", 1),
            ("조직 변화와 실제 워크플로 통합", 1),
            "",
            "빠른 실행이 승패를 가름:",
            ("조직 재설계", 1),
            ("업무 프로세스 재구성", 1),
            ("직원 스킬 변화", 1),
            ("검증된 외부 솔루션 적극 활용", 1),
            "",
            "완전 자율이 아닌 인간-AI 협업 모델이 현실적 해답"
        ])
    ]

    # 6페이지(인덱스 5) 이후에 슬라이드 삽입
    insert_slides_after_index(prs, slides_data, 5)

    # 파일 저장
    prs.save('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')
    print(f"AI Agent 도입 실패 슬라이드가 성공적으로 추가되었습니다!")
    print(f"총 {len(prs.slides)}개의 슬라이드가 포함되어 있습니다.")
    print(f"{len(slides_data)}개의 슬라이드가 6페이지 이후에 삽입되었습니다.")

if __name__ == "__main__":
    main()
