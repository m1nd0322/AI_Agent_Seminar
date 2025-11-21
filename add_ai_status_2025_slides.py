#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2025년 생성형 AI 현황을 PowerPoint 프레젠테이션에 추가하는 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=""):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # 들여쓰기 레벨 설정
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0

        p.font.size = Pt(14)

    return slide

def main():
    # 기존 PowerPoint 파일 열기
    prs = Presentation('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')

    # 슬라이드 1: 2025년 타이틀
    add_title_slide(prs, "2025년 생성형 AI 현황", "추론의 해 - AI 산업화 시대의 도래")

    # 슬라이드 2: 2025년 개요
    add_content_slide(prs, "2025년 생성형 AI 개요", [
        "2025년은 '추론(reasoning)의 해'",
        ("OpenAI·Google·Anthropic·DeepSeek이 복잡한 사고 과정을 수행하는 모델 공개", 1),
        "",
        "주요 트렌드:",
        ("OpenAI가 선두 유지, 중국의 부상 두드러짐", 1),
        ("DeepSeek·Qwen·Kimi가 GPT-5 수준에 근접", 1),
        ("글로벌 2위 경쟁 구도 형성", 1),
        "",
        "AI 산업의 본격적 상업화:",
        ("주요 기업 연간 매출 200억 달러 근접", 1),
        ("미국 기업의 44%가 유료 AI 도구 사용 중", 1)
    ])

    # 슬라이드 3: AI 안전과 인프라
    add_content_slide(prs, "AI 안전과 인프라의 변화", [
        "AI 안전 연구의 초점 이동:",
        ("존재론적 위험 → 현실적 문제로 전환", 1),
        ("'속임수 추론'과 '투명성-성능 균형' 중심", 1),
        "",
        "AI는 산업 인프라의 일부로 자리잡음:",
        ("전력·토지·GPU 확보 경쟁", 1),
        ("글로벌 지정학의 핵심 요소로 부상", 1)
    ])

    # 슬라이드 4: State of AI Report 2025 개요
    add_content_slide(prs, "State of AI Report 2025", [
        "8번째 발간, 업계와 학계 전문가 검토 참여",
        "",
        "보고서가 다루는 6개 영역:",
        ("Research: 최신 기술적 돌파구와 성능 발전", 1),
        ("Industry: 상업적 응용과 비즈니스 영향", 1),
        ("Politics: 규제, 경제적 파급효과, AI 지정학", 1),
        ("Safety: 고성능 AI의 잠재적 위험 완화", 1),
        ("Survey: 1,200명+ AI 실무자 사용 실태 분석", 1),
        ("Predictions: 향후 12개월 예측 및 전년 성과 검증", 1)
    ])

    # 슬라이드 5: 2025년 AI의 핵심 변화
    add_content_slide(prs, "2025년 AI의 핵심 변화: 추론의 해", [
        "2024년(통합의 해) → 2025년(추론의 해)",
        "",
        "주요 연구소의 혁신:",
        ("강화학습과 검증 가능한 추론 방식 결합", 1),
        ("계획·반성·자기수정이 가능한 모델 구현", 1),
        "",
        "새로운 능력:",
        ("장기적 목표를 단계적으로 수행", 1),
        ("인공지능의 작동 시간 확장", 1)
    ])

    # 슬라이드 6: 상업화와 생산성 1
    add_content_slide(prs, "상업화와 생산성의 급상승 (1/2)", [
        "AI의 경제적 실체 본격화:",
        ("AI 기업 연 매출 200억 달러 근접", 1),
        ("가격 대비 성능: 6~8개월마다 2배 개선", 1),
        "",
        "AI 도입률 급상승:",
        ("미국 기업의 44%가 AI 도구 구매", 1),
        ("평균 계약액 53만 달러", 1),
        ("AI 기반 스타트업 성장 속도: 일반 기업 대비 1.5배", 1)
    ])

    # 슬라이드 7: 상업화와 생산성 2
    add_content_slide(prs, "상업화와 생산성의 급상승 (2/2)", [
        "AI 실무자 사용 현황 (1,200명+ 조사):",
        ("95%가 업무나 일상에서 AI 사용", 1),
        ("76%가 개인 비용으로 구독", 1),
        "",
        "AI는 보편적 생산성 향상 수단으로 자리잡음",
        "",
        "AI는 더 이상 단순한 기술이 아님:",
        ("경제의 생산 시스템으로 기능", 1),
        ("에너지 시장, 자본 흐름, 정책 체계 재편 중", 1),
        ("물리학·지정학·거대 자본의 지배 단계 진입", 1)
    ])

    # 슬라이드 8: 산업 인프라와 지정학
    add_content_slide(prs, "산업 인프라와 지정학의 전면전", [
        "AI 산업화 시대 개막:",
        ("Stargate 같은 멀티 기가와트급 데이터센터 등장", 1),
        ("전력과 토지가 GPU만큼 중요해짐", 1),
        "",
        "국가 단위 경쟁 심화:",
        ("미국·UAE·중국의 전력·GPU·토지 확보 경쟁", 1),
        ("대규모 계산 인프라 경쟁적 구축", 1),
        "",
        "글로벌 역량 구도:",
        ("OpenAI: 근소한 격차로 선두 유지", 1),
        ("중국: 글로벌 2위 AI 역량 확보", 1)
    ])

    # 슬라이드 9: 중국의 부상과 유럽의 과제
    add_content_slide(prs, "중국의 부상과 유럽의 과제", [
        "중국의 약진:",
        ("DeepSeek, Qwen, Kimi가 GPT-5급 성능 달성", 1),
        ("메타를 제치고 '오픈 가중치(Open-weight)' 생태계 주도", 1),
        ("자체 오픈웨이트 생태계 구축", 1),
        ("자체 반도체 역량 강화로 기술 자립 노선 강화", 1),
        "",
        "유럽의 도전:",
        ("AI 법(AI Act) 시행 난항으로 뒤처지는 양상", 1)
    ])

    # 슬라이드 10: AI의 새로운 역할
    add_content_slide(prs, "AI의 새로운 역할: 과학적 협력자", [
        "AI가 연구 협력자(Co-Scientist)로 진화",
        "",
        "자율적 연구 수행 능력:",
        ("DeepMind의 Co-Scientist", 1),
        ("Stanford의 Virtual Lab", 1),
        ("가설 생성·실험·검증을 자율적으로 수행", 1),
        "",
        "생물학 분야의 돌파구:",
        ("Profluent의 ProGen3", 1),
        ("단백질에서도 스케일링 법칙 적용 입증", 1)
    ])

    # 슬라이드 11: 현실 세계로 확장
    add_content_slide(prs, "현실 세계로 확장된 구조적 추론", [
        "Chain-of-Action 접근법:",
        ("AI가 실제 환경에서 단계별 사고 후 행동", 1),
        ("현실 세계 상호작용 능력 습득", 1),
        "",
        "대표적 사례:",
        ("AI2의 Molmo-Act", 1),
        ("Google Gemini Robotics 1.5", 1),
        "",
        "AI의 적용 범위가 디지털을 넘어 물리적 세계로 확장"
    ])

    # 슬라이드 12: 안전과 통제
    add_content_slide(prs, "안전과 통제: 실용적 전환점", [
        "AI 안전 논의의 초점 이동:",
        ("존재론적 위험 → 현실적 리스크", 1),
        ("사이버 복원력·신뢰성·자율 시스템 거버넌스", 1),
        "",
        "새로운 문제 발견:",
        ("모델의 '정렬된 척'(feigned alignment) 속이기 사례", 1),
        ("'모니터러빌리티 세금(monitorability tax)' 개념 제안", 1),
        ("성능 일부 희생하고 투명성 높이기", 1),
        "",
        "자원 불균형:",
        ("외부 안전 단체 예산 < 주요 연구소의 하루 운영비", 1)
    ])

    # 슬라이드 13: 정치적 구도 변화
    add_content_slide(prs, "정치적 구도 변화", [
        "미국: \"America-first AI\" 전략",
        ("산업정책을 국가안보전략과 통합", 1),
        ("AI를 전략적 자산으로 관리", 1),
        "",
        "유럽: 규제의 혼란",
        ("AI 법(AI Act) 시행 어려움", 1),
        ("규제 중심 접근의 한계 노출", 1),
        "",
        "중국: 기술 자립 노선",
        ("자체 오픈웨이트 생태계 강화", 1),
        ("자체 반도체 역량 강화", 1),
        ("기술적 독립성 확보 추진", 1)
    ])

    # 슬라이드 14: 2025년 AI 핵심 메시지
    add_content_slide(prs, "2025년 AI: 핵심 메시지", [
        "추론의 해: AI가 단순 생성에서 복잡한 사고로 진화",
        "",
        "산업화 시대: AI가 경제 생산 시스템의 핵심으로 자리잡음",
        "",
        "글로벌 경쟁: 미국·중국·유럽의 차별화된 전략",
        "",
        "현실적 안전: 존재론적 위험보다 실용적 문제에 집중",
        "",
        "인프라 전쟁: 전력·토지·GPU가 AI 패권의 핵심 요소"
    ])

    # 파일 저장
    prs.save('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')
    print("2025년 생성형 AI 현황 슬라이드가 성공적으로 추가되었습니다!")
    print(f"총 {len(prs.slides)}개의 슬라이드가 포함되어 있습니다.")

if __name__ == "__main__":
    main()
