#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint 파일의 슬라이드 정보를 확인하는 스크립트
"""

from pptx import Presentation

def main():
    # 기존 PowerPoint 파일 열기
    prs = Presentation('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')

    print(f"총 슬라이드 수: {len(prs.slides)}")
    print("\n슬라이드 정보:")

    for i, slide in enumerate(prs.slides):
        print(f"\n슬라이드 {i+1} (인덱스 {i}):")
        print(f"  레이아웃: {slide.slide_layout.name}")

        # 제목이 있는 경우
        if slide.shapes.title:
            print(f"  제목: {slide.shapes.title.text}")

        # 6페이지만 상세히 출력
        if i == 5:
            print("\n  === 6페이지 상세 정보 ===")
            for shape_idx, shape in enumerate(slide.shapes):
                print(f"    Shape {shape_idx}: {shape.shape_type}")
                if hasattr(shape, "text"):
                    print(f"      텍스트: {shape.text[:100] if len(shape.text) > 100 else shape.text}")

if __name__ == "__main__":
    main()
