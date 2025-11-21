#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2024년 생성형 AI 현황을 PowerPoint 프레젠테이션에 추가하는 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=""):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # 들여쓰기 레벨 설정
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0

        p.font.size = Pt(14)

    return slide

def main():
    # 기존 PowerPoint 파일 열기
    prs = Presentation('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')

    # 슬라이드 1: 타이틀
    add_title_slide(prs, "2024년 생성형 AI 현황", "파일럿에서 프로덕션으로의 전환")

    # 슬라이드 2: AI 도입 가속화
    add_content_slide(prs, "파일럿에서 프로덕션으로: AI 도입 가속화", [
        "2024년은 생성형 AI가 기업의 핵심 전략으로 자리 잡은 해",
        "",
        "주요 지표:",
        ("AI 관련 지출: $23억(2023) → $138억(2024), 6배 이상 증가", 1),
        ("기업의 72%가 생성형 AI 도구 채택 확대 예상", 1),
        "",
        "많은 조직이 구현 전략을 명확히 정의하지 못했지만,",
        "대규모 변화의 초기 단계로 진입 중"
    ])

    # 슬라이드 3: 생성형 AI 지출
    add_content_slide(prs, "생성형 AI 지출: 기업의 헌신적 투자", [
        "예산 조달 방식:",
        ("60%: 혁신 예산에서 조달", 1),
        ("40%: 기존 예산 재배치 (AI를 장기 전략으로 편입)", 1),
        "",
        "$46억이 생성형 AI 애플리케이션에 투자",
        ("전년 대비 8배 증가", 1),
        "",
        "기업들은 평균 10개의 잠재적 활용 사례를 식별하며",
        "도입 속도 증가"
    ])

    # 슬라이드 4: 주요 사용 사례와 ROI
    add_content_slide(prs, "주요 사용 사례와 ROI", [
        "코드 코파일럿 (51% 도입)",
        ("GitHub Copilot, Codeium, Cursor 등으로 개발 생산성 향상", 1),
        "",
        "지원 챗봇 (31% 도입)",
        ("24/7 지원 제공: Aisera, Decagon, Sierra 등", 1),
        "",
        "엔터프라이즈 검색 및 데이터 변환 (28%, 27% 도입)",
        ("데이터 사일로 통합: Glean, Sana 등", 1),
        "",
        "회의 요약 (24% 도입)",
        ("Fireflies.ai, Otter.ai, Fathom으로 회의 효율 개선", 1)
    ])

    # 슬라이드 5: AI 기반 에이전트와 자동화
    add_content_slide(prs, "AI 기반 에이전트와 자동화의 부상", [
        "현재:",
        ("인간의 작업을 보완하는 방식 선호", 1),
        "",
        "미래:",
        ("독립적으로 복잡한 프로세스를 관리하는", 1),
        ("완전 자동화 에이전트로 전환 예상", 1),
        "",
        "대표 사례: Clay, Forge"
    ])

    # 슬라이드 6: 도입 장벽과 실패 원인
    add_content_slide(prs, "도입 장벽과 실패 원인", [
        "주요 실패 이유:",
        ("구현 비용 (26%)", 1),
        ("데이터 프라이버시 문제 (21%)", 1),
        ("기대에 못 미치는 ROI (18%)", 1),
        ("기술적 문제 (15%, \"헛소리\" 생성 포함)", 1),
        "",
        "성공을 위한 고려사항:",
        ("ROI와 맞춤화 우선", 1),
        ("기술 통합 및 지원 체계 구축", 1)
    ])

    # 슬라이드 7: 산업별 생성형 AI 도입 현황
    add_content_slide(prs, "산업별 생성형 AI 도입 현황", [
        "헬스케어 ($5억 지출)",
        ("자동화된 의료 문서화 및 환자 관리: Abridge, Notable", 1),
        "",
        "법률 ($3.5억 지출)",
        ("계약 검토 및 소송 준비 자동화: Harvey, Everlaw", 1),
        "",
        "금융 서비스 ($1억 지출)",
        ("회계 및 규제 준수 개선: Numeric, Rogo", 1),
        "",
        "미디어 및 엔터테인먼트 ($1억 지출)",
        ("콘텐츠 제작 효율 향상: Runway, Descript", 1)
    ])

    # 슬라이드 8: AI 기술 스택과 설계 패턴
    add_content_slide(prs, "AI 기술 스택과 설계 패턴", [
        "주요 트렌드:",
        ("RAG(검색 증강 생성): 51% 채택", 1),
        ("미세 조정: 9%만 도입 (여전히 드물)", 1),
        ("에이전틱 아키텍처: 12%에서 도입 시작", 1),
        "",
        "데이터 스토리지:",
        ("AI 네이티브 벡터 DB Pinecone 18% 시장 점유", 1)
    ])

    # 슬라이드 9: 전망 개요
    add_content_slide(prs, "전망: AI의 미래에 대한 주요 예측", [
        "2024년은 하이프에서 현실적 구현으로 전환된 해",
        "",
        "향후 세 가지 주요 예측:",
        "",
        ("1. AI 에이전트가 차세대 변화를 주도", 1),
        ("2. \"다윗이 골리앗을 이긴다\": 기존 기업의 쇠퇴", 1),
        ("3. AI 인재 부족 심화", 1)
    ])

    # 슬라이드 10: 전망 1 - AI 에이전트
    add_content_slide(prs, "전망 1: AI 에이전트가 차세대 변화를 주도", [
        "에이전틱 자동화가 AI 변혁을 주도:",
        ("콘텐츠 생성과 지식 검색을 넘어", 1),
        ("복잡한 다단계 작업을 처리", 1),
        "",
        "시장 영향:",
        ("Clay와 Forge 같은 플랫폼이", 1),
        ("$4천억 소프트웨어 시장과", 1),
        ("$10조 미국 서비스 경제에 도전", 1),
        "",
        "새로운 인프라 필요:",
        ("에이전트 인증, 도구 통합 플랫폼,", 1),
        ("AI 브라우저 프레임워크, AI 생성 코드 전용 런타임 등", 1)
    ])

    # 슬라이드 11: 전망 2 - 기존 기업의 쇠퇴
    add_content_slide(prs, "전망 2: \"다윗이 골리앗을 이긴다\" - 기존 기업의 쇠퇴", [
        "ChatGPT의 영향:",
        ("Chegg: 시장 가치 85% 소멸", 1),
        ("Stack Overflow: 웹 트래픽 절반 감소", 1),
        "",
        "향후 도전이 예상되는 영역:",
        ("IT 아웃소싱 기업 (Cognizant)", 1),
        ("기존 자동화 기업 (UiPath)", 1),
        ("대형 소프트웨어 기업 (Salesforce, Autodesk)", 1),
        ("→ AI 네이티브 스타트업의 도전 가능성", 1)
    ])

    # 슬라이드 12: 전망 3 - AI 인재 부족
    add_content_slide(prs, "전망 3: AI 인재 부족 심화", [
        "AI 시스템 확산과 고도화로 심각한 인재 부족 초래",
        "",
        "필요한 인재:",
        ("데이터 과학자뿐 아니라", 1),
        ("AI 기술과 도메인 전문성을 결합할 수 있는 전문가", 1),
        "",
        "경쟁 심화:",
        ("AI 숙련 엔터프라이즈 아키텍트 연봉", 1),
        ("2-3배 상승 현상 보편화", 1),
        "",
        "훈련 프로그램 및 AI 센터 투자에도 불구,",
        "인재 수요는 공급을 초과",
        "",
        "다음 AI 혁신을 이끌 인재 확보를 위한 치열한 경쟁 예상"
    ])

    # 파일 저장
    prs.save('/home/user/AI_Agent_Seminar/AI_Agent_Seminar.pptx')
    print("2024년 생성형 AI 현황 슬라이드가 성공적으로 추가되었습니다!")
    print(f"총 {len(prs.slides)}개의 슬라이드가 포함되어 있습니다.")

if __name__ == "__main__":
    main()
